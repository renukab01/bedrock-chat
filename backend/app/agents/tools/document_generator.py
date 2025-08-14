import logging
import json
import io
import base64
import re
import os
import uuid
from typing import Any, Dict, List, Optional
from datetime import datetime

import boto3
from app.agents.tools.agent_tool import AgentTool
from app.repositories.models.custom_bot import BotModel
from app.routes.schemas.conversation import type_model_name
from app.utils import generate_presigned_url, BEDROCK_REGION
from pydantic import BaseModel, Field

# Document generation libraries
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# S3 bucket for document storage - using environment variable or default
# The bucket name should be set by the CDK deployment
DOCUMENT_BUCKET = os.environ.get("DOCUMENT_BUCKET", f"bedrock-chat-documents-{os.environ.get('ENV_NAME', 'default')}")


def _upload_document_to_s3(document_data: bytes, filename: str, content_type: str) -> str:
    """
    Upload document to S3 and return presigned download URL.
    
    Args:
        document_data: Binary document data
        filename: Name of the file
        content_type: MIME type of the document
        
    Returns:
        str: Presigned download URL
    """
    try:
        # Log bucket information for debugging
        logger.info(f"Attempting to upload to bucket: {DOCUMENT_BUCKET}")
        logger.info(f"Using region: {BEDROCK_REGION}")
        logger.info(f"ENV_NAME: {os.environ.get('ENV_NAME', 'not_set')}")
        logger.info(f"DOCUMENT_BUCKET env var: {os.environ.get('DOCUMENT_BUCKET', 'not_set')}")
        
        # Generate unique S3 key with timestamp and UUID to avoid conflicts
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]
        s3_key = f"generated_documents/{timestamp}_{unique_id}_{filename}"
        
        # Create S3 client
        s3_client = boto3.client("s3", region_name=BEDROCK_REGION)
        
        # First, try to check if bucket exists and is accessible
        try:
            s3_client.head_bucket(Bucket=DOCUMENT_BUCKET)
            logger.info(f"Bucket {DOCUMENT_BUCKET} exists and is accessible")
        except Exception as bucket_error:
            logger.error(f"Bucket check failed: {bucket_error}")
            # List available buckets for debugging
            try:
                buckets = s3_client.list_buckets()
                bucket_names = [bucket['Name'] for bucket in buckets['Buckets']]
                logger.info(f"Available buckets: {bucket_names}")
            except Exception as list_error:
                logger.error(f"Could not list buckets: {list_error}")
            raise bucket_error
        
        # Upload to S3
        s3_client.put_object(
            Bucket=DOCUMENT_BUCKET,
            Key=s3_key,
            Body=document_data,
            ContentType=content_type,
            # Set expiration for cleanup (30 days)
            Expires=datetime.now().timestamp() + (30 * 24 * 60 * 60)
        )
        
        # Generate presigned download URL (valid for 24 hours)
        download_url = generate_presigned_url(
            bucket=DOCUMENT_BUCKET,
            key=s3_key,
            expiration=24 * 60 * 60,  # 24 hours
            client_method="get_object"
        )
        
        logger.info(f"Document uploaded successfully to S3: {s3_key}")
        return download_url
        
    except Exception as e:
        logger.error(f"Error uploading document to S3: {e}")
        logger.error(f"Bucket: {DOCUMENT_BUCKET}, Region: {BEDROCK_REGION}")
        raise e


def _sanitize_filename(filename: str) -> str:
    """
    Sanitize filename to comply with AWS Bedrock document restrictions.
    This matches the _convert_to_valid_file_name function in conversation.py:
    - Only alphanumeric characters, whitespace characters, hyphens, parentheses, and square brackets
    - No more than one consecutive whitespace character
    """
    # Remove invalid characters (keep only alphanumeric, spaces, hyphens, parentheses, square brackets)
    sanitized = re.sub(r"[^a-zA-Z0-9\s\-\(\)\[\]]", "", filename)
    
    # Replace multiple consecutive whitespace characters with single space
    sanitized = re.sub(r"\s+", " ", sanitized)
    
    # Trim whitespace from start and end
    sanitized = sanitized.strip()
    
    # If the filename is empty after sanitization, provide a default
    if not sanitized:
        sanitized = "document"
    
    return sanitized


class ExcelGeneratorInput(BaseModel):
    title: str = Field(description="Title for the Excel document")
    data: List[Dict[str, Any]] = Field(description="Data to populate the Excel sheet. Each dict represents a row with column names as keys.")
    sheet_name: str = Field(default="Sheet1", description="Name of the Excel sheet")
    include_header: bool = Field(default=True, description="Whether to include column headers")


class WordGeneratorInput(BaseModel):
    title: str = Field(description="Title for the Word document")
    content: List[Dict[str, Any]] = Field(description="Content sections. Each dict should have 'type' (heading, paragraph, list) and 'text' or 'items'")


class PowerPointGeneratorInput(BaseModel):
    title: str = Field(description="Title for the PowerPoint presentation")
    slides: List[Dict[str, Any]] = Field(description="Slide content. Each dict should have 'title' and 'content' (list of bullet points or paragraphs)")


def _generate_excel(tool_input: ExcelGeneratorInput, bot: BotModel | None, model: type_model_name | None) -> Dict[str, Any]:
    """Generate an Excel file from the provided data."""
    try:
        logger.info(f"Generating Excel document: {tool_input.title}")
        
        # Create workbook and worksheet
        wb = Workbook()
        ws = wb.active
        ws.title = tool_input.sheet_name
        
        # Add title
        ws['A1'] = tool_input.title
        ws['A1'].font = Font(bold=True, size=16)
        ws.merge_cells('A1:E1')
        ws['A1'].alignment = Alignment(horizontal='center')
        
        # Add data
        if tool_input.data:
            start_row = 3
            
            # Add headers if requested
            if tool_input.include_header and tool_input.data:
                headers = list(tool_input.data[0].keys())
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=start_row, column=col, value=header)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
                start_row += 1
            
            # Add data rows
            for row_idx, row_data in enumerate(tool_input.data, start_row):
                for col_idx, (key, value) in enumerate(row_data.items(), 1):
                    ws.cell(row=row_idx, column=col_idx, value=str(value))
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = None
            for cell in column:
                try:
                    # Skip merged cells that don't have column_letter attribute
                    if hasattr(cell, 'column_letter'):
                        column_letter = cell.column_letter
                        if cell.value and len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                except:
                    pass
            
            if column_letter:
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        # Save to bytes
        excel_buffer = io.BytesIO()
        wb.save(excel_buffer)
        excel_buffer.seek(0)
        
        # Create filename with proper sanitization
        sanitized_title = _sanitize_filename(tool_input.title)
        filename = f"{sanitized_title}.xlsx"
        
        # Upload to S3 and get download URL
        download_url = _upload_document_to_s3(
            document_data=excel_buffer.getvalue(),
            filename=filename,
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        
        logger.info(f"Generated presigned URL for Excel: {download_url}")
        
        # Return result with download link
        result = {
            "content": f"Excel document '{sanitized_title}' has been generated successfully with {len(tool_input.data)} rows of data. The file contains columns: {', '.join(tool_input.data[0].keys()) if tool_input.data else 'No data'}. Download URL: {download_url}",
            "source_name": filename,
            "source_link": download_url
        }
        
        logger.info(f"Returning Excel result: {result}")
        return result
        
    except Exception as e:
        logger.error(f"Error generating Excel document: {e}")
        raise e


def _generate_word(tool_input: WordGeneratorInput, bot: BotModel | None, model: type_model_name | None) -> Dict[str, Any]:
    """Generate a Word document from the provided content."""
    try:
        logger.info(f"Generating Word document: {tool_input.title}")
        
        # Create document
        doc = Document()
        
        # Add title
        title_paragraph = doc.add_heading(tool_input.title, level=1)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add timestamp
        doc.add_paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph()  # Empty line
        
        # Add content sections
        for section in tool_input.content:
            try:
                # Handle both dict and string inputs
                if isinstance(section, str):
                    # If it's a string, treat it as a paragraph
                    doc.add_paragraph(section)
                    continue
                    
                if not isinstance(section, dict):
                    # If it's neither string nor dict, convert to string and add as paragraph
                    doc.add_paragraph(str(section))
                    continue
                
                section_type = section.get('type', 'paragraph')
                
                if section_type == 'heading':
                    level = section.get('level', 2)
                    text = section.get('text', '')
                    if text:
                        doc.add_heading(text, level=level)
                    
                elif section_type == 'paragraph':
                    text = section.get('text', '')
                    if text:
                        doc.add_paragraph(text)
                    
                elif section_type == 'list':
                    items = section.get('items', [])
                    if items:
                        for item in items:
                            doc.add_paragraph(str(item), style='List Bullet')
                        
                elif section_type == 'table':
                    table_data = section.get('data', [])
                    if table_data and len(table_data) > 0:
                        # Ensure all rows have the same number of columns
                        max_cols = max(len(row) if isinstance(row, list) else 1 for row in table_data)
                        table = doc.add_table(rows=len(table_data), cols=max_cols)
                        table.style = 'Table Grid'
                        
                        for row_idx, row_data in enumerate(table_data):
                            if isinstance(row_data, list):
                                for col_idx, cell_data in enumerate(row_data):
                                    if col_idx < max_cols:
                                        table.cell(row_idx, col_idx).text = str(cell_data)
                            else:
                                table.cell(row_idx, 0).text = str(row_data)
                else:
                    # Unknown section type, treat as paragraph
                    text = section.get('text', str(section))
                    doc.add_paragraph(text)
                    
            except Exception as section_error:
                logger.warning(f"Error processing section {section}: {section_error}")
                # Fallback: add the section as a paragraph
                doc.add_paragraph(str(section))
        
        # Save to bytes
        word_buffer = io.BytesIO()
        doc.save(word_buffer)
        word_buffer.seek(0)
        
        # Create filename with proper sanitization
        sanitized_title = _sanitize_filename(tool_input.title)
        filename = f"{sanitized_title}.docx"
        
        # Upload to S3 and get download URL
        download_url = _upload_document_to_s3(
            document_data=word_buffer.getvalue(),
            filename=filename,
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
        logger.info(f"Generated presigned URL for Word: {download_url}")
        
        # Return result with download link
        result = {
            "content": f"Word document '{sanitized_title}' has been generated successfully with {len(tool_input.content)} content sections including headings, paragraphs, and lists. Download URL: {download_url}",
            "source_name": filename,
            "source_link": download_url
        }
        
        logger.info(f"Returning Word result: {result}")
        return result
        
    except Exception as e:
        logger.error(f"Error generating Word document: {e}")
        raise e


def _generate_powerpoint(tool_input: PowerPointGeneratorInput, bot: BotModel | None, model: type_model_name | None) -> Dict[str, Any]:
    """Generate a PowerPoint presentation using python-pptx."""
    try:
        logger.info(f"Generating PowerPoint presentation: {tool_input.title}")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = PptxInches(13.33)
        prs.slide_height = PptxInches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        # Set title and subtitle
        title_slide.shapes.title.text = tool_input.title
        if title_slide.shapes.placeholders[1]:  # Subtitle placeholder
            title_slide.shapes.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add content slides
        for slide_data in tool_input.slides:
            slide_title = slide_data.get('title', 'Slide')
            content = slide_data.get('content', [])
            
            # Use bullet slide layout
            bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Set slide title
            slide.shapes.title.text = slide_title
            
            # Add content to the content placeholder
            if slide.shapes.placeholders[1]:  # Content placeholder
                content_placeholder = slide.shapes.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear default text
                
                # Add bullet points
                for i, item in enumerate(content):
                    if i == 0:
                        # First paragraph (already exists)
                        p = text_frame.paragraphs[0]
                    else:
                        # Add new paragraphs for subsequent items
                        p = text_frame.add_paragraph()
                    
                    p.text = str(item)
                    p.level = 0  # Top level bullet
                    p.alignment = PP_ALIGN.LEFT
        
        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Create filename with proper sanitization
        sanitized_title = _sanitize_filename(tool_input.title)
        filename = f"{sanitized_title}.pptx"
        
        # Upload to S3 and get download URL
        download_url = _upload_document_to_s3(
            document_data=pptx_buffer.getvalue(),
            filename=filename,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
        logger.info(f"Generated presigned URL for PowerPoint: {download_url}")
        
        # Return result with download link
        result = {
            "content": f"PowerPoint presentation '{sanitized_title}' has been generated successfully with {len(tool_input.slides) + 1} slides (including title slide). The presentation includes a title slide and {len(tool_input.slides)} content slides with bullet points. Download URL: {download_url}",
            "source_name": filename,
            "source_link": download_url
        }
        
        logger.info(f"Returning PowerPoint result: {result}")
        return result
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint presentation: {e}")
        raise e


# Create the agent tools
excel_generator_tool = AgentTool(
    name="excel_generator",
    description="Generate an Excel spreadsheet with data. Useful for creating reports, data tables, and structured information.",
    args_schema=ExcelGeneratorInput,
    function=_generate_excel,
)

word_generator_tool = AgentTool(
    name="word_generator", 
    description="Generate a Word document with formatted content. Useful for creating reports, documentation, and text-based documents.",
    args_schema=WordGeneratorInput,
    function=_generate_word,
)

powerpoint_generator_tool = AgentTool(
    name="powerpoint_generator",
    description="Generate a PowerPoint presentation with slides. Useful for creating presentations, slide decks, and visual content.",
    args_schema=PowerPointGeneratorInput,
    function=_generate_powerpoint,
)